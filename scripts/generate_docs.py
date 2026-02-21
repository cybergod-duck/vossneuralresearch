import os
from docx import Document
from pptx import Presentation
from openpyxl import Workbook

output_dir = os.path.join(os.path.dirname(__file__), '..', 'public', 'documents')
os.makedirs(output_dir, exist_ok=True)

# 1. Clinical Presentation.pptx
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Voss Neural Research LLC"
slide.placeholders[1].text = "Empirical Study of AI in Stimulant Use Disorder Recovery\nN=1 Clinical Case Study"

bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
slide2.shapes.title.text = "The 4 Voss Protocols"
tf = slide2.shapes.placeholders[1].text_frame
tf.text = "Protocol I: Defined Exits"
tf.add_paragraph().text = "Protocol II: Somatic Anchoring"
tf.add_paragraph().text = "Protocol III: Output Ownership"
tf.add_paragraph().text = "Protocol IV: The Bridge"

prs.save(os.path.join(output_dir, 'Clinical_Presentation.pptx'))
print("Created Clinical_Presentation.pptx")

# 2. Provider Handout.docx
doc1 = Document()
doc1.add_heading('Provider Handout: AI & SUD Recovery', 0)
doc1.add_paragraph('What to know when your patient is using AI in recovery.')
doc1.add_heading('Red Flags to Watch For:', level=1)
doc1.add_paragraph('Complete absence of physical output bridging', style='List Bullet')
doc1.add_paragraph('Rapidly diminishing sleep schedules', style='List Bullet')
doc1.add_paragraph('Extreme hostility when AI access is restricted', style='List Bullet')
doc1.add_heading('Boundary Recommendations:', level=1)
doc1.add_paragraph('Ensure the patient utilizes the 4 Voss Protocols to prevent behavior transference. Ask to review their Bridge Checklist and weekly commitment card.')
doc1.save(os.path.join(output_dir, 'Provider_Handout.docx'))
print("Created Provider_Handout.docx")

# 3. Extended Research Report.docx
doc2 = Document()
doc2.add_heading('Extended Research Report', 0)
doc2.add_heading('Comparative Analysis: AI vs Behavioral Addictions', level=1)
doc2.add_paragraph('While gaming and social media hijack passive consumption pathways, excessive AI generation loops ("hyper-ideation") hijack active execution reward pathways. The dopamine trigger is linked to instantaneous realization of high-level conceptual ideas without the friction of execution.')
doc2.add_heading('Clinical Trial Design Proposal', level=1)
doc2.add_paragraph('Future N=20 longitudinal study utilizing the Voss Neural Recovery Metrics tracker.')
doc2.save(os.path.join(output_dir, 'Extended_Research_Report.docx'))
print("Created Extended_Research_Report.docx")

# 4. Recovery Metrics.xlsx
wb = Workbook()
ws = wb.active
ws.title = "Daily Logs"
ws.append(["Date", "AI Building (hrs)", "AI Being (hrs)", "Sleep (Quality 1-10)", "Cravings (1-10)", "Bridge Checklist Met?", "Offline Output Score"])
ws.append(["2025-01-01", 6.5, 2.0, 7, 3, "Yes", 85])
ws.append(["2025-01-02", 8.0, 1.5, 6, 5, "No", 40])
wb.save(os.path.join(output_dir, 'Recovery_Metrics.xlsx'))
print("Created Recovery_Metrics.xlsx")
